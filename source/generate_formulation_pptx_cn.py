from pptx import Presentation
from pptx.util import Inches, Pt
import pptx

def create_presentation():
    # Create presentation with 16:9 aspect ratio (widescreen)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "问题建模：效用与隐私的张力"
    subtitle.text = "偏斜感知图流摘要中的隐私泄露形式化分析\n\nP-Scube 的理论基础"

    # Slide 2: The Context
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "1. 背景定义：图流与偏斜感知"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "图流与幂律分布 (Power-Law Distribution)："
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• 真实世界的图流呈现极度的结构偏斜。"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 节点 v 的度数服从幂律分布：Pr[d_v = k] ∝ k^(-γ)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "偏斜感知资源分配 (以 Scube 为例)："
    p.font.bold = True
    p.space_before = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• 为了最小化哈希冲突，Scube 根据估计度数 d̂_v 为节点分配资源 A_v (哈希地址数)。"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 分配函数是一个确定性的阶跃函数 (Deterministic Step Function)："
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "  A_v = f(d̂_v) = max(2, ⌈d̂_v / θ⌉)"
    p.level = 1
    p.font.italic = True
    
    p = tf.add_paragraph()
    p.text = "核心洞察：函数 f(·) 是严格单调递增且完全确定性的。这虽然带来了极高的查询效用，但也埋下了严重的隐私隐患。"
    p.level = 1
    p.font.bold = True
    p.font.color.rgb = pptx.dml.color.RGBColor(200, 0, 0)

    # Slide 3: The Adversary Model
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "2. 攻击者模型 (The Adversary Model)"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "攻击者的观测能力 (Observation)："
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• 节点查询的延迟 T_v 与系统为其分配的地址数 A_v 呈严格的线性正相关。"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "  T_v ≈ c · A_v + noise_network"
    p.level = 1
    p.font.italic = True
    
    p = tf.add_paragraph()
    p.text = "• 结论：攻击者仅需通过查询接口 (Query-only API) 测量时间侧信道，即可高精度推断出 A_v。"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "攻击者的目标 (Goal - Hub Extraction)："
    p.font.bold = True
    p.space_before = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• 判断目标节点 v 是否为“高度节点 (Hub)” (即真实度数 d_v > τ)。"
    p.level = 1

    # Slide 4: Formalizing Privacy Leakage
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "3. 隐私泄露的形式化量化"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "先验概率 (Prior Probability)："
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• 在未观测到 A_v 时，由于幂律分布，随机盲猜一个节点是 Hub 的概率极低。"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "  Pr[d_v > τ] ≈ 0.01"
    p.level = 1
    p.font.italic = True
    
    p = tf.add_paragraph()
    p.text = "后验概率 (Posterior Probability)："
    p.font.bold = True
    p.space_before = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• 攻击者观测到系统分配的地址数 A_v 后，猜中 Hub 的概率变为 Pr[d_v > τ | A_v]。"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "可提取性优势 (Extractability Advantage, Adv)："
    p.font.bold = True
    p.space_before = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• 定义为后验概率与先验概率的差值："
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "  Adv(A) = | Pr[d_v > τ | A_v] - Pr[d_v > τ] |"
    p.level = 1
    p.font.italic = True
    
    p = tf.add_paragraph()
    p.text = "• Adv(A) 越接近 1，说明隐私泄露越严重 (攻击者获得了近乎完美的提取能力)。"
    p.level = 1

    # Slide 5: The Tension Theorem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "4. 效用-隐私张力定理 (The Tension Theorem)"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "基于确定性映射的贝叶斯推断："
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "  Pr[d_v > τ | A_v] = (Pr[A_v | d_v > τ] · Pr[d_v > τ]) / Pr[A_v]"
    p.level = 1
    p.font.italic = True
    
    p = tf.add_paragraph()
    p.text = "• 因为分配机制 A_v = ⌈d̂_v / θ⌉ 是完全确定性的，如果观测到 A_v 很大："
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "  Pr[A_v | d_v > τ] ≈ 1  且  Pr[A_v | d_v < τ] ≈ 0"
    p.level = 2
    p.font.italic = True
    
    p = tf.add_paragraph()
    p.text = "定理推导结果："
    p.font.bold = True
    p.space_before = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• 条件概率的坍缩导致后验概率无限逼近于 100%。"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "  Adv(A) ≈ 1 - Pr[d_v > τ] ≈ 0.99"
    p.level = 1
    p.font.italic = True
    
    p = tf.add_paragraph()
    p.text = "结论：在图流摘要中，确定性的偏斜感知资源分配 (Skew-awareness) 在数学上必然导致攻击者获得近乎完美的 Hub 提取能力。这就是不可调和的效用-隐私张力。"
    p.level = 1
    p.font.bold = True
    p.font.color.rgb = pptx.dml.color.RGBColor(200, 0, 0)

    # Slide 6: The Mitigation Strategy (P-Scube)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "5. 破局思路：P-Scube 的防御机制"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "核心目标：打破确定性，引入度数不可区分性 (Degree Indistinguishability)"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• 必须在数学上限制攻击者区分两个相似度数节点的能力 (满足 Local DP)："
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "  Pr[A_u* = k] / Pr[A_v* = k] ≤ e^ε"
    p.level = 1
    p.font.italic = True
    
    p = tf.add_paragraph()
    p.text = "P-Scube 的双重机制："
    p.font.bold = True
    p.space_before = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "1. 单调加噪升级 (Monotonic Noisy Promotion)："
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "  d̃_v = d̂_v + Laplace(λ)"
    p.level = 2
    p.font.italic = True
    
    p = tf.add_paragraph()
    p.text = "  (度数为 500 的节点也有概率得分超过 600 的升级阈值，打破确定性)"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "2. 粗粒度资源分配 (Coarsened Allocation)："
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "  将细粒度的 A_v ∈ {2,3,4,5,6} 映射为粗粒度桶 A_v* ∈ {2, 4, 6}"
    p.level = 2
    p.font.italic = True
    
    p = tf.add_paragraph()
    p.text = "  (度数 400 和 600 的节点被强行塞入同一个“4地址桶”，彻底阻断细粒度推断)"
    p.level = 2

    prs.save('/Volumes/external_disk_macos /thesis/Scube/source/problem_formulation_cn.pptx')

import pptx
create_presentation()