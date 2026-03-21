from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Problem Formulation: The Utility-Privacy Tension"
    subtitle.text = "Formalizing Privacy Leakage in Skew-Aware Graph Summarization\n\nTheoretical Foundations for P-Scube"

    # Slide 2: The Context (Graph Streams & Skew)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "1. The Context: Graph Streams & Skew"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "Graph Streams and Power-Law Distribution:"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Real-world graph streams exhibit structural skew."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• The degree of node v follows a power-law: Pr[d_v = k] ∝ k^(-γ)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Skew-Aware Allocation (e.g., Scube):"
    p.font.bold = True
    p.space_before = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• To minimize hash collisions, Scube allocates resources A_v (addresses) based on estimated degree d̂_v."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Allocation is a deterministic step function: A_v = f(d̂_v) = max(2, ⌈d̂_v / θ⌉)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Insight: The function f(·) is strictly monotonically increasing and completely deterministic. This provides high utility but introduces severe privacy risks."
    p.level = 1
    p.font.italic = True
    p.font.color.rgb = pptx.dml.color.RGBColor(200, 0, 0)

    # Slide 3: The Adversary Model
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "2. The Adversary Model"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "Adversary's Observation (A_v):"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Query latency T_v is linearly correlated with allocated addresses A_v."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• T_v ≈ c · A_v + noise_network"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Adversary can accurately infer A_v via timing side-channels (Query-only access)."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Adversary's Goal (Hub Extraction):"
    p.font.bold = True
    p.space_before = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• Determine if a target node v is a 'Hub' (i.e., true degree d_v > τ)."
    p.level = 1

    # Slide 4: Formalizing Privacy Leakage
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "3. Formalizing Privacy Leakage"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "Prior Probability:"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Without observing A_v, the probability of guessing a hub is extremely low (due to power-law)."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Pr[d_v > τ] ≈ 0.01"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Posterior Probability:"
    p.font.bold = True
    p.space_before = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• After observing A_v, the probability becomes Pr[d_v > τ | A_v]."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Extractability Advantage (Adv(A)):"
    p.font.bold = True
    p.space_before = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• Adv(A) = | Pr[d_v > τ | A_v] - Pr[d_v > τ] |"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Adv(A) → 1 means perfect extraction (severe leakage)."
    p.level = 1

    # Slide 5: The Tension Theorem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "4. The Utility-Privacy Tension Theorem"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "Bayesian Inference on Deterministic Mapping:"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Pr[d_v > τ | A_v] = (Pr[A_v | d_v > τ] · Pr[d_v > τ]) / Pr[A_v]"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Because A_v = ⌈d̂_v / θ⌉ is deterministic, if A_v is large, Pr[A_v | d_v > τ] ≈ 1 and Pr[A_v | d_v < τ] ≈ 0."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Theorem Result:"
    p.font.bold = True
    p.space_before = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• The posterior probability collapses to 100% certainty."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Adv(A) ≈ 1 - Pr[d_v > τ] ≈ 0.99"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Conclusion: Skew-aware resource allocation mathematically guarantees near-perfect hub extraction. This is the fundamental Utility-Privacy Tension."
    p.level = 1
    p.font.bold = True
    p.font.color.rgb = pptx.dml.color.RGBColor(200, 0, 0)

    # Slide 6: The Mitigation Strategy (P-Scube)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "5. The Mitigation Strategy: P-Scube"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "Goal: Break the Deterministic Mapping"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• We must enforce Degree Indistinguishability (Local DP constraint):"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Pr[A_u* = k] / Pr[A_v* = k] ≤ e^ε  (for nodes u, v with similar degrees)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "P-Scube Mechanisms:"
    p.font.bold = True
    p.space_before = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "1. Monotonic Noisy Promotion: d̂_v → d̃_v = d̂_v + Laplace(λ)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • A node with degree 500 now has a probability to surpass the 600 threshold."
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "2. Coarsened Allocation: Map fine-grained A_v ∈ {2,3,4,5,6} to coarse buckets A_v* ∈ {2, 4, 6}."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Nodes with degrees 400 and 600 are forced into the same '4-address bucket', destroying the adversary's ability to distinguish them."
    p.level = 2

    prs.save('/Volumes/external_disk_macos /thesis/Scube/source/problem_formulation.pptx')

import pptx
create_presentation()
