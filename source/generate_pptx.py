import sys
sys.path.insert(0, "/Volumes/external_disk_macos /thesis/Scube/.pip_packages")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG       = RGBColor(0x1B, 0x1B, 0x2F)
CARD     = RGBColor(0x27, 0x27, 0x4A)
ACCENT   = RGBColor(0x6C, 0x63, 0xFF)
CYAN     = RGBColor(0x00, 0xD2, 0xFF)
RED      = RGBColor(0xFF, 0x6B, 0x6B)
GREEN    = RGBColor(0x00, 0xE6, 0x96)
ORANGE   = RGBColor(0xFF, 0xA0, 0x40)
YELLOW   = RGBColor(0xFF, 0xE0, 0x66)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0xBB, 0xBB, 0xCC)
DIM      = RGBColor(0x88, 0x88, 0xAA)


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def box(slide, l, t, w, h, fill, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb


def bullets(slide, l, t, w, h, items, sz=18, color=GRAY, dot_color=CYAN, spacing=Pt(12)):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        r1 = p.add_run()
        r1.text = "●  "
        r1.font.size = Pt(sz)
        r1.font.color.rgb = dot_color
        r1.font.name = "Calibri"
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(sz)
        r2.font.color.rgb = color
        r2.font.name = "Calibri"


def line(slide, l, t, w):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(3))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()


def tag(slide, text, l=Inches(0.6), t=Inches(0.45)):
    s = box(slide, l, t, Inches(2.2), Inches(0.38), ACCENT)
    tf = s.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER


# ══════════════════════════════════════
# 1 — Title
# ══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
line(s, Inches(1.8), Inches(2.2), Inches(9.8))
txt(s, Inches(1.4), Inches(2.55), Inches(10.6), Inches(1.2),
    "P-Scube：偏斜感知图流摘要中的结构隐私泄露", sz=34, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1.5), Inches(4.0), Inches(10.3), Inches(0.8),
    "Skew-aware optimization improves hub utility, but also amplifies hub privacy leakage.",
    sz=20, color=GRAY, align=PP_ALIGN.CENTER)
line(s, Inches(1.8), Inches(5.0), Inches(9.8))
txt(s, Inches(1.6), Inches(5.45), Inches(10.2), Inches(0.5),
    "主线：high-degree leakage -> selective mitigation -> tradeoff evaluation",
    sz=16, color=DIM, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════
# 2 — Background & Tension
# ══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tag(s, "BACKGROUND")
txt(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
    "研究背景：为什么 Scube 值得讨论隐私？", sz=30, bold=True)
txt(s, Inches(0.6), Inches(1.75), Inches(12), Inches(0.45),
    "Scube 的强点，恰恰可能是隐私风险的来源。", sz=18, color=GRAY)

steps = [
    ("1", "普通 sketch", "所有节点近似同等对待", GRAY),
    ("2", "SCube", "优先保障 hub", CYAN),
    ("3", "副作用", "hub 更容易被识别", RED),
]
for i, (num, title, desc, clr) in enumerate(steps):
    x = Inches(0.7 + i * 4.15)
    box(s, x, Inches(2.7), Inches(3.6), Inches(2.1), CARD, clr)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.35), Inches(2.15), Inches(0.8), Inches(0.8))
    circ.fill.solid()
    circ.fill.fore_color.rgb = clr
    circ.line.fill.background()
    p = circ.text_frame.paragraphs[0]
    p.text = num
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BG
    p.alignment = PP_ALIGN.CENTER
    txt(s, x + Inches(0.25), Inches(3.15), Inches(3.1), Inches(0.4), title, sz=22, bold=True, color=clr, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.25), Inches(3.75), Inches(3.1), Inches(0.45), desc, sz=16, color=GRAY, align=PP_ALIGN.CENTER)

box(s, Inches(0.8), Inches(5.45), Inches(11.8), Inches(1.0), CARD, ACCENT)
txt(s, Inches(1.05), Inches(5.72), Inches(11.2), Inches(0.3),
    "一句话：hub utility gain 可能伴随 hub privacy leakage gain。", sz=20, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════
# 3 — SCube Core Idea
# ══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tag(s, "SCUBE IDEA")
txt(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
    "在讨论隐私前，先看 Scube 的中心思想", sz=30, bold=True)
txt(s, Inches(0.6), Inches(1.7), Inches(11.6), Inches(0.5),
    "Scube 的本质不是“更大的 sketch”，而是“有选择地优待高度节点”。", sz=18, color=GRAY)

core_parts = [
    ("① Detect", "DegDetector 估计节点活跃度", "谁更活跃，谁更可能成为 hub", CYAN),
    ("② Allocate", "给 hub 更多地址", "普通节点 2 个地址，hub 可以有 3、4、5...", GREEN),
    ("③ Benefit", "减少冲突、提高查询质量", "因此 hub 的 accuracy 与 latency 都更好", ORANGE),
]
for i, (title, sub, desc, clr) in enumerate(core_parts):
    x = Inches(0.55 + i * 4.2)
    box(s, x, Inches(2.6), Inches(3.8), Inches(2.2), CARD, clr)
    txt(s, x + Inches(0.2), Inches(2.85), Inches(3.4), Inches(0.35), title, sz=24, bold=True, color=clr, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.2), Inches(3.35), Inches(3.4), Inches(0.35), sub, sz=16, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.2), Inches(4.0), Inches(3.4), Inches(0.65), desc, sz=15, color=GRAY, align=PP_ALIGN.CENTER)

box(s, Inches(0.75), Inches(5.65), Inches(11.8), Inches(0.95), CARD, ACCENT)
txt(s, Inches(1.05), Inches(5.9), Inches(11.2), Inches(0.3),
    "关键理解：Scube 通过“差异化服务”换来性能提升，而隐私问题正是从这种差异化开始的。",
    sz=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════
# 4 — Focus & Threat Model
# ══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tag(s, "FOCUS")
txt(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
    "论文主线：问题收缩而不是全面铺开", sz=30, bold=True)

box(s, Inches(0.75), Inches(2.0), Inches(5.4), Inches(1.75), CARD, CYAN)
txt(s, Inches(1.05), Inches(2.22), Inches(4.8), Inches(0.35), "只回答两个问题", sz=21, bold=True, color=CYAN)
txt(s, Inches(1.05), Inches(2.75), Inches(4.7), Inches(0.7),
    "1. `SCube` 是否更易泄露 hub？\n2. 能否低代价缓解？", sz=18, color=WHITE)

task_titles = [("主攻击 1", "Top-K Hub"), ("主攻击 2", "Hub Membership"), ("辅攻击", "Timing")]
task_colors = [RED, GREEN, YELLOW]
for i, ((k, v), clr) in enumerate(zip(task_titles, task_colors)):
    x = Inches(6.55 + i * 2.05)
    box(s, x, Inches(2.0), Inches(1.9), Inches(1.75), CARD, clr)
    txt(s, x + Inches(0.1), Inches(2.35), Inches(1.7), Inches(0.25), k, sz=12, bold=True, color=clr, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.1), Inches(2.8), Inches(1.7), Inches(0.35), v, sz=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

box(s, Inches(0.75), Inches(4.4), Inches(11.9), Inches(1.55), CARD, ORANGE)
txt(s, Inches(1.05), Inches(4.65), Inches(11.2), Inches(0.3), "答辩时不展开", sz=20, bold=True, color=ORANGE)
txt(s, Inches(1.05), Inches(5.1), Inches(11.0), Inches(0.45),
    "`edge probing / occupancy / kick-out / ext / 全局DP` 统统降级，否则故事线会散。",
    sz=18, color=GRAY)


# ══════════════════════════════════════
# 5 — Demo Evidence I
# ══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tag(s, "PRIVACY DEMO")
txt(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
    "Privacy Demo：实验设置与 Hub 提取结果", sz=30, bold=True)

box(s, Inches(0.75), Inches(1.95), Inches(4.2), Inches(3.55), CARD, CYAN)
txt(s, Inches(1.05), Inches(2.15), Inches(3.5), Inches(0.35), "实验设置", sz=22, bold=True, color=CYAN)
bullets(s, Inches(1.05), Inches(2.65), Inches(3.45), Inches(2.2), [
    "5000 个节点，903,858 条边",
    "Zipf 度分布，`gamma = 1.2`",
    "矩阵大小 `3000 x 3000`",
    "真实度统计：Median = 19，P99 = 1978，Max = 2261",
], sz=15, color=GRAY)

box(s, Inches(5.15), Inches(1.95), Inches(7.45), Inches(3.55), CARD, RED)
txt(s, Inches(5.45), Inches(2.15), Inches(6.8), Inches(0.35), "Attack 1: Top-K Hub Extraction", sz=22, bold=True, color=RED)
bullets(s, Inches(5.45), Inches(2.65), Inches(6.7), Inches(2.1), [
    "对全部节点执行 `nodeWeightQuery(v, 0)` 并排序",
    "`Precision@10 = 100%`，`Precision@20 = 100%`",
    "`Precision@50 = 98%`，`Precision@100 = 100%`",
    "真实 top-50 hub 的平均攻击排名 = `25.52`，最佳 = 1，最差 = 51",
    "Hub 与 Non-hub 的平均相对误差都接近 `0%`",
], sz=15, color=GRAY)

box(s, Inches(0.85), Inches(5.8), Inches(11.7), Inches(0.85), CARD, ACCENT)
txt(s, Inches(1.15), Inches(6.03), Inches(11.1), Inches(0.3),
    "结论 1：仅凭公开节点查询，攻击者几乎可以直接恢复 hub 排名。", sz=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════
# 6 — Demo Evidence II
# ══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tag(s, "PRIVACY DEMO")
txt(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
    "Privacy Demo：边探测与时延侧信道", sz=30, bold=True)

box(s, Inches(0.7), Inches(1.95), Inches(5.8), Inches(3.7), CARD, ORANGE)
txt(s, Inches(1.0), Inches(2.15), Inches(5.1), Inches(0.35), "Attack 2: Sensitive Edge Probing", sz=22, bold=True, color=ORANGE)
bullets(s, Inches(1.0), Inches(2.65), Inches(5.0), Inches(1.8), [
    "测试集：100 条真实边 + 100 条不存在边",
    "结果：`TP = 100`，`FP = 0`，`TN = 100`，`FN = 0`",
    "`Precision = 100%`，`Recall = 100%`，`FPR = 0%`",
    "样例中既包含 hub，也包含普通节点，判断仍然完全正确",
], sz=15, color=GRAY)

box(s, Inches(6.75), Inches(1.95), Inches(5.85), Inches(3.7), CARD, YELLOW)
txt(s, Inches(7.05), Inches(2.15), Inches(5.15), Inches(0.35), "Attack 3: Timing Side-Channel", sz=22, bold=True, color=YELLOW)
bullets(s, Inches(7.05), Inches(2.65), Inches(5.0), Inches(1.8), [
    "Hub 平均延迟 = `4.973 us`，Non-hub = `3.404 us`",
    "Latency ratio = `1.4608x`",
    "分类结果：`Accuracy = 92.64%`，`Recall = 100%`",
    "Pearson(`degree`, `latency`) = `0.5037`",
], sz=15, color=GRAY)

box(s, Inches(0.85), Inches(5.9), Inches(11.7), Inches(0.8), CARD, ACCENT)
txt(s, Inches(1.15), Inches(6.12), Inches(11.1), Inches(0.3),
    "结论 2：即使对查询值加噪，查询路径和耗时仍然会泄露 degree class。", sz=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════
# 7 — Root Cause
# ══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tag(s, "ROOT CAUSE")
txt(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
    "泄露根源：hub-oriented optimization 本身暴露了 hub", sz=30, bold=True)
flow = [
    ("更多地址", CYAN),
    ("更少冲突", GREEN),
    ("更高 utility", ORANGE),
    ("更易识别 hub", RED),
]
for i, (label, clr) in enumerate(flow):
    x = Inches(0.7 + i * 3.1)
    box(s, x, Inches(2.9), Inches(2.3), Inches(1.3), CARD, clr)
    txt(s, x + Inches(0.15), Inches(3.25), Inches(2.6), Inches(0.3), label, sz=20, bold=True, color=clr, align=PP_ALIGN.CENTER)

box(s, Inches(0.8), Inches(4.8), Inches(5.75), Inches(1.2), CARD, CYAN)
txt(s, Inches(1.1), Inches(5.05), Inches(5.1), Inches(0.3), "Signal A", sz=14, bold=True, color=CYAN)
txt(s, Inches(1.1), Inches(5.4), Inches(5.1), Inches(0.35), "Detector state\n暴露 degree level", sz=18, color=WHITE)

box(s, Inches(6.75), Inches(4.8), Inches(5.15), Inches(1.2), CARD, GREEN)
txt(s, Inches(7.05), Inches(5.05), Inches(4.5), Inches(0.3), "Signal B", sz=14, bold=True, color=GREEN)
txt(s, Inches(7.05), Inches(5.4), Inches(4.5), Inches(0.35), "Query path\n暴露 value + latency", sz=18, color=WHITE)

txt(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.35),
    "一句话：utility gain 与 leakage gain 来自同一套 skew-aware 机制。", sz=18, color=RED, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════
# 8 — P-Scube Design
# ══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tag(s, "P-SCUBE")
txt(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
    "方法设计：P-Scube 只保留两个核心机制", sz=30, bold=True)
txt(s, Inches(0.6), Inches(1.7), Inches(11.5), Inches(0.5),
    "不去掉 skew-awareness，而是打破 `degree -> address_number` 的精确映射。", sz=18, color=CYAN, bold=True)

mechanisms = [
    ("机制 1", "Noisy Degree Promotion", CYAN,
     "把“确定升级”改成“概率升级”",
     "重点：削弱 top-K 排名与 address level 的稳定映射",
     "代价：少量 hub 可能延迟升级",
     "实现点：promotion threshold 加入随机扰动"),
    ("机制 2", "Coarsened Address Allocation", GREEN,
     "把 2,3,4,5... 粗化成 2 / 4 / 6",
     "重点：多个节点共享同一 observable class",
     "代价：部分节点会拿到略多地址",
     "实现点：只修改 detector 返回的地址级别"),
]
for i, (mid, title, clr, line1, line2, line3, line4) in enumerate(mechanisms):
    x = Inches(0.65 + i * 6.2)
    box(s, x, Inches(2.35), Inches(5.95), Inches(3.1), CARD, clr)
    txt(s, x + Inches(0.25), Inches(2.6), Inches(5.35), Inches(0.25), mid, sz=13, bold=True, color=clr)
    txt(s, x + Inches(0.25), Inches(2.95), Inches(5.35), Inches(0.4), title, sz=23, bold=True, color=clr)
    txt(s, x + Inches(0.25), Inches(3.5), Inches(5.2), Inches(0.35), line1, sz=17, color=WHITE)
    txt(s, x + Inches(0.25), Inches(4.15), Inches(5.2), Inches(0.45), "注释 1: " + line2, sz=15, color=GRAY)
    txt(s, x + Inches(0.25), Inches(4.85), Inches(5.2), Inches(0.45), "注释 2: " + line3, sz=15, color=GRAY)
    txt(s, x + Inches(0.25), Inches(5.55), Inches(5.2), Inches(0.45), "注释 3: " + line4, sz=15, color=GRAY)

box(s, Inches(0.95), Inches(6.55), Inches(11.55), Inches(0.55), CARD, ACCENT)
txt(s, Inches(1.2), Inches(6.7), Inches(11.0), Inches(0.25),
    "设计原则：只改 detector 侧，不大改矩阵主体，这样更容易保留 Scube 的性能优势。", sz=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════
# 9 — Unified Metrics
# ══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tag(s, "METRICS")
txt(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
    "统一评价框架：把攻击、缓解和代价放到同一语言里", sz=30, bold=True)
box(s, Inches(0.9), Inches(2.0), Inches(11.3), Inches(1.2), CARD, CYAN)
txt(s, Inches(1.2), Inches(2.28), Inches(10.7), Inches(0.3),
    "AttackAdv = AttackSuccess - RandomBaseline", sz=24, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

eqs = [
    ("LeakageGain", "`AttackAdv(SCube) - AttackAdv(GSS)`", ORANGE),
    ("MitigationGain", "`AttackAdv(SCube) - AttackAdv(P-SCube)`", GREEN),
]
for i, (title, expr, clr) in enumerate(eqs):
    x = Inches(1.2 + i * 5.8)
    box(s, x, Inches(3.8), Inches(4.9), Inches(1.3), CARD, clr)
    txt(s, x + Inches(0.2), Inches(4.05), Inches(4.5), Inches(0.25), title, sz=18, bold=True, color=clr, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.2), Inches(4.45), Inches(4.5), Inches(0.3), expr, sz=16, color=WHITE, align=PP_ALIGN.CENTER)

box(s, Inches(1.1), Inches(5.55), Inches(10.8), Inches(0.9), CARD, ACCENT)
txt(s, Inches(1.35), Inches(5.83), Inches(10.3), Inches(0.28),
    "再配合 `ARE` 和 latency overhead，就得到最关键的一张图：privacy-utility frontier。",
    sz=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════
# 10 — Evaluation Plan
# ══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tag(s, "EVALUATION")
txt(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
    "实验设计：围绕一张 tradeoff 主图组织", sz=30, bold=True)
box(s, Inches(0.8), Inches(2.0), Inches(2.8), Inches(1.4), CARD, CYAN)
txt(s, Inches(1.0), Inches(2.25), Inches(2.4), Inches(0.3), "Baselines", sz=22, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
txt(s, Inches(1.0), Inches(2.72), Inches(2.4), Inches(0.4), "`GSS`\n`SCube`\n`P-SCube`", sz=17, color=WHITE, align=PP_ALIGN.CENTER)

box(s, Inches(3.95), Inches(2.0), Inches(2.8), Inches(1.4), CARD, ORANGE)
txt(s, Inches(4.2), Inches(2.25), Inches(2.3), Inches(0.3), "Main Tasks", sz=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
txt(s, Inches(4.2), Inches(2.72), Inches(2.3), Inches(0.4), "Hub extraction\nMembership\nTiming", sz=17, color=WHITE, align=PP_ALIGN.CENTER)

box(s, Inches(7.1), Inches(2.0), Inches(5.0), Inches(1.4), CARD, GREEN)
txt(s, Inches(7.35), Inches(2.25), Inches(4.5), Inches(0.3), "Must-Prove Claims", sz=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
txt(s, Inches(7.35), Inches(2.72), Inches(4.5), Inches(0.4), "SCube > GSS in leakage\nP-Scube < SCube at low utility cost", sz=17, color=WHITE, align=PP_ALIGN.CENTER)

box(s, Inches(1.0), Inches(4.15), Inches(11.3), Inches(1.9), CARD, ACCENT)
txt(s, Inches(1.35), Inches(4.45), Inches(10.6), Inches(0.32), "答辩主图", sz=22, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
txt(s, Inches(1.35), Inches(4.95), Inches(10.6), Inches(0.35),
    "x 轴：protection strength", sz=18, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, Inches(1.35), Inches(5.35), Inches(10.6), Inches(0.35),
    "y 轴：AttackAdv    vs    ARE / latency overhead", sz=18, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, Inches(1.0), Inches(6.35), Inches(11.4), Inches(0.28),
    "一句话：不是找最好参数，而是证明 P-Scube 的 frontier 更优。", sz=18, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════
# 11 — Contributions & Close
# ══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tag(s, "CONTRIBUTIONS")
txt(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
    "三点贡献与总结", sz=30, bold=True)

contribs = [
    ("问题", "发现 hub utility 与 hub leakage 的耦合", CYAN),
    ("方法", "`NoisyPromotion` + `CoarsenedAllocation`", GREEN),
    ("评估", "`AttackAdv / LeakageGain / frontier`", ORANGE),
]
for i, (title, desc, clr) in enumerate(contribs):
    x = Inches(0.8 + i * 4.05)
    box(s, x, Inches(2.4), Inches(3.6), Inches(2.4), CARD, clr)
    txt(s, x + Inches(0.2), Inches(2.75), Inches(3.2), Inches(0.35), title, sz=24, bold=True, color=clr, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.2), Inches(3.45), Inches(3.2), Inches(0.75), desc, sz=17, color=WHITE, align=PP_ALIGN.CENTER)

box(s, Inches(0.9), Inches(5.9), Inches(11.8), Inches(0.8), CARD, RED)
txt(s, Inches(1.2), Inches(6.08), Inches(11.2), Inches(0.3),
    "一句话 takeaway：Better hub utility in SCube leads to better hub extractability; P-Scube weakens this coupling.",
    sz=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ── Save ──
out = "/Volumes/external_disk_macos /thesis/Scube/P-Scube_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
